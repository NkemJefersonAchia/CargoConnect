#!/usr/bin/env python3
"""CargoConnect – 10-slide visual presentation (7-minute deck)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0A, 0x1F, 0x44)
ORANGE   = RGBColor(0xE8, 0x5D, 0x04)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xF4, 0xF6, 0xFB)
TEXT     = RGBColor(0x33, 0x41, 0x55)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
BORDER   = RGBColor(0xDD, 0xE3, 0xEF)
GREEN    = RGBColor(0x06, 0x5F, 0x46)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
AMBER    = RGBColor(0x92, 0x40, 0x0E)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
DARK     = RGBColor(0x1E, 0x2A, 0x3A)
NAVY2    = RGBColor(0x1A, 0x3A, 0x6E)
RED_BG   = RGBColor(0xFF, 0xE4, 0xE6)
RED      = RGBColor(0x9F, 0x12, 0x39)
BLUE_BG  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE     = RGBColor(0x1E, 0x40, 0xAF)

W, H = 13.333, 7.5
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────
def ns(bg_color=WHITE):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg_color
    return s

def box(s, l, t, w, h, fc=None, lc=None, lw=0.75, rr=False):
    sh = s.shapes.add_shape(5 if rr else 1,
                            Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fc else sh.fill.background()
    if fc: sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:  sh.line.fill.background()
    return sh

def circ(s, l, t, d, fc):
    sh = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, sz=13, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.name = 'Calibri'
    r.font.color.rgb = color if color else TEXT
    return bx

def mtf(s, l, t, w, h):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    return tf

def ap(tf, text, sz=11, bold=False, color=None, align=PP_ALIGN.LEFT,
       before=0, after=3, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    p.space_after = Pt(after)
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.name = 'Calibri'
    r.font.color.rgb = color if color else TEXT

# ── Layout helpers ────────────────────────────────────────────────────────────
def hdr(s, title, num, light=False):
    box(s, 0, 0, W, 0.96, fc=NAVY)
    box(s, 0, 0.92, W, 0.06, fc=ORANGE)
    txt(s, f"{num:02d}", 0.35, 0.1, 0.55, 0.78, sz=10, bold=True, color=ORANGE)
    txt(s, title, 0.82, 0.1, 12.1, 0.78, sz=23, bold=True, color=WHITE)

def ftr(s, label=""):
    box(s, 0, H - 0.28, W, 0.28, fc=NAVY)
    txt(s, "CargoConnect", 0.35, H - 0.26, 4, 0.22, sz=8, color=MUTED)
    if label:
        txt(s, label, W - 5.0, H - 0.26, 4.7, 0.22,
            sz=8, color=MUTED, align=PP_ALIGN.RIGHT)

def num_circ(s, l, t, d, n, fc=ORANGE):
    circ(s, l, t, d, fc)
    txt(s, str(n), l, t + d * 0.15, d, d * 0.72,
        sz=int(d * 14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def stat_card(s, l, t, w, h, num, label):
    box(s, l, t, w, h, fc=NAVY, rr=True)
    txt(s, num,   l, t + 0.1,       w, h * 0.58,
        sz=44, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, label, l, t + h * 0.6,   w, h * 0.36,
        sz=10, color=MUTED, align=PP_ALIGN.CENTER)

def divider(s, l, t, w):
    box(s, l, t, w, 0.03, fc=BORDER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 01 – TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = ns(NAVY)
box(s, 0, 0,        W, 0.16, fc=ORANGE)
box(s, 0, H - 0.16, W, 0.16, fc=ORANGE)
box(s, 0.65, 1.3, 0.08, 4.1, fc=ORANGE)            # left accent bar

txt(s, "CargoConnect", 0.95, 1.25, 9, 1.9,
    sz=72, bold=True, color=WHITE)
txt(s, "Smart Logistics for Kigali, Rwanda",
    0.95, 3.12, 8.5, 0.6, sz=18,
    color=RGBColor(0xBF, 0xCA, 0xE0))
box(s, 0.95, 3.82, 3.2, 0.05, fc=ORANGE)
txt(s, "TECHNICAL PRESENTATION  |  TRIMESTER 4",
    0.95, 3.94, 7, 0.32, sz=10, bold=True, color=ORANGE)

# Team cards (right side)
box(s, 9.3, 1.35, 3.6, 5.0, fc=NAVY2, rr=True)
txt(s, "TEAM", 9.3, 1.52, 3.6, 0.3,
    sz=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
for i, name in enumerate(["[Member 1]", "[Member 2]",
                           "[Member 3]", "[Member 4]"]):
    box(s, 9.5, 2.02 + i * 0.9, 3.2, 0.66, fc=NAVY, rr=True)
    txt(s, name, 9.5, 2.06 + i * 0.9, 3.2, 0.6,
        sz=13, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "March 2026", 9.3, 6.1, 3.6, 0.26,
    sz=9, color=MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 02 – PROBLEM + SOLUTION
# ═════════════════════════════════════════════════════════════════════════════
s = ns(NAVY)
box(s, 0, 0, W, 0.16, fc=ORANGE)
box(s, 0, H - 0.16, W, 0.16, fc=ORANGE)

# ── LEFT: problem cards ───────────────────────────────────────────────────────
txt(s, "The Problem", 0.3, 0.3, 5.5, 0.7,
    sz=22, bold=True, color=WHITE)

pain = [
    ("No Marketplace",     "Drivers found by word of mouth"),
    ("No Transparency",    "Arbitrary, unverifiable pricing"),
    ("No Live Tracking",   "Zero visibility during transit"),
    ("No Digital Payment", "Cash-only, no mobile money"),
]
for i, (title, sub) in enumerate(pain):
    y = 1.12 + i * 1.47
    box(s, 0.3, y, 5.7, 1.32, fc=NAVY2, rr=True)
    circ(s, 0.44, y + 0.42, 0.48, fc=ORANGE)
    txt(s, "!", 0.44, y + 0.44, 0.48, 0.38,
        sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.04, y + 0.2,  4.8, 0.42, sz=13, bold=True, color=WHITE)
    txt(s, sub,   1.04, y + 0.65, 4.8, 0.38, sz=10,
        color=RGBColor(0xBF, 0xCA, 0xE0))

# ── DIVIDER ───────────────────────────────────────────────────────────────────
box(s, 6.3, 0.3, 0.05, H - 0.62, fc=RGBColor(0x2A, 0x4A, 0x7A))

# ── RIGHT: solution diagram ────────────────────────────────────────────────────
txt(s, "The Solution", 6.6, 0.3, 6.5, 0.7,
    sz=22, bold=True, color=WHITE)

# Platform hub
box(s, 8.35, 3.0, 2.9, 1.3, fc=ORANGE, rr=True)
txt(s, "CargoConnect", 8.35, 3.08, 2.9, 0.65,
    sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Web Platform", 8.35, 3.62, 2.9, 0.34,
    sz=9, color=WHITE, align=PP_ALIGN.CENTER)

# Three roles
role_positions = [
    ("CUSTOMER",  "Book  |  Track  |  Pay",  6.6,  3.1),
    ("DRIVER",    "Accept  |  Navigate  |  Complete",  11.1, 3.1),
    ("ADMIN",     "Verify  |  Monitor  |  Manage",     8.5,  5.35),
]
for label, sub, rx, ry in role_positions:
    box(s, rx, ry, 2.36, 1.18, fc=NAVY, rr=True)
    txt(s, label, rx, ry + 0.08, 2.36, 0.42,
        sz=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, sub,   rx, ry + 0.54, 2.36, 0.36,
        sz=8.5, color=RGBColor(0xBF, 0xCA, 0xE0), align=PP_ALIGN.CENTER)

# Arrows
txt(s, ">", 9.0,  3.52, 0.3, 0.3, sz=14, bold=True, color=ORANGE)   # left->hub
txt(s, "<", 11.1 - 0.35, 3.52, 0.3, 0.3, sz=14, bold=True, color=ORANGE) # hub->right
box(s, 9.8, 4.32, 0.05, 1.05, fc=ORANGE)                              # hub->admin
txt(s, "v", 9.76, 5.28, 0.16, 0.22, sz=12, bold=True, color=ORANGE)

ftr(s, "Problem and Solution")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 03 – PLATFORM AT A GLANCE + FEATURES
# ═════════════════════════════════════════════════════════════════════════════
s = ns(BG)
hdr(s, "Platform at a Glance", 3)

# Stat row
stats = [("44","API Routes"), ("7","Blueprints"), ("8","DB Models"),
         ("14","Templates"),  ("3","User Roles"), ("1","Payment API")]
sw = 2.06
for i, (num, lbl) in enumerate(stats):
    stat_card(s, 0.32 + i*(sw + 0.06), 1.06, sw, 1.72, num, lbl)

# Feature grid (done badges)
done = ["User Auth", "RBAC", "Driver Search", "Cost Estimation",
        "Booking Lifecycle", "Live Tracking", "MoMo Payments",
        "Notifications", "Admin Panel", "Rating System"]
next_ = ["Email Alerts", "Surge Pricing", "PWA"]

txt(s, "IMPLEMENTED", 0.32, 2.98, 3, 0.28, sz=8, bold=True, color=GREEN)
x, y = 0.32, 3.3
for f in done:
    w2 = len(f) * 0.082 + 0.35
    if x + w2 > W - 0.3:
        x = 0.32; y += 0.42
    box(s, x, y, w2, 0.32, fc=GREEN_BG, lc=RGBColor(0x6E,0xE7,0xB7), lw=0.6, rr=True)
    txt(s, f, x + 0.08, y + 0.04, w2 - 0.12, 0.26,
        sz=9, bold=True, color=GREEN)
    x += w2 + 0.12

y += 0.52
txt(s, "NEXT PHASE", 0.32, y, 2.5, 0.28, sz=8, bold=True, color=AMBER)
y += 0.3; x = 0.32
for f in next_:
    w2 = len(f) * 0.082 + 0.35
    box(s, x, y, w2, 0.32, fc=AMBER_BG, lc=RGBColor(0xFB,0xBF,0x24), lw=0.6, rr=True)
    txt(s, f, x + 0.08, y + 0.04, w2 - 0.12, 0.26,
        sz=9, bold=True, color=AMBER)
    x += w2 + 0.12

ftr(s, "Features Overview")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 04 – TECHNOLOGY STACK
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "Technology Stack", 4)

rows = [
    ("Backend",       "Flask 3.x",            "vs Django / FastAPI",
     "Lightweight  |  Application Factory  |  Blueprints"),
    ("Database",      "PostgreSQL",            "vs MySQL / SQLite",
     "Native ENUMs  |  ACID compliance  |  Concurrent writes"),
    ("Real-Time",     "Flask-SocketIO",        "vs HTTP Polling",
     "Bidirectional WebSocket  |  Per-booking rooms"),
    ("Maps",          "Leaflet.js",            "vs Google Maps",
     "Open-source  |  No API cost  |  Rwanda tile support"),
    ("Payments",      "MTN MoMo API",          "vs Stripe",
     "Rwanda-dominant  |  USSD mobile money  |  RWF currency"),
    ("Auth",          "Flask-Login + bcrypt",  "vs JWT",
     "Server sessions  |  Adaptive hash cost"),
]

# Column header
box(s, 0.3, 1.05, W - 0.6, 0.38, fc=NAVY)
for label, cx, cw in [("Area", 0.4, 1.7), ("Chosen Tool", 2.14, 2.8),
                       ("Alternative", 5.0, 2.7), ("Why", 7.76, 5.25)]:
    txt(s, label, cx, 1.07, cw, 0.32, sz=9, bold=True, color=WHITE)

rh = 0.83
for ri, (area, chosen, alt, why) in enumerate(rows):
    ry = 1.45 + ri * rh
    rc = BG if ri % 2 == 0 else WHITE
    box(s, 0.3, ry, W - 0.6, rh - 0.04, fc=rc, lc=BORDER, lw=0.3)
    # Area chip
    box(s, 0.3, ry, 1.8, rh - 0.04, fc=NAVY)
    txt(s, area, 0.34, ry + 0.22, 1.72, 0.4,
        sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Chosen (green pill)
    box(s, 2.14, ry + 0.16, 2.76, 0.5, fc=GREEN_BG,
        lc=RGBColor(0x6E,0xE7,0xB7), lw=0.6, rr=True)
    txt(s, chosen, 2.22, ry + 0.22, 2.6, 0.38,
        sz=10, bold=True, color=GREEN)
    # Alternative (muted)
    txt(s, alt, 5.06, ry + 0.24, 2.6, 0.38,
        sz=9, color=MUTED, italic=True)
    # Why
    txt(s, why, 7.82, ry + 0.22, 5.1, 0.42, sz=10, bold=True, color=TEXT)

ftr(s, "Technology Stack")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 05 – ALGORITHM + DATABASE
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "Algorithm and Database Design", 5)

# ── LEFT: Algorithm visual ────────────────────────────────────────────────────
# Pipeline steps
steps = [("Filter", "Verified +\nAvailable"),
         ("Capacity", "Truck >=\ncargo"),
         ("Distance", "Haversine\nformula"),
         ("Rank", "Top 10\nby rating")]
sw2 = 1.48; sy = 1.08
for i, (title, sub) in enumerate(steps):
    x = 0.3 + i * (sw2 + 0.22)
    bg2 = ORANGE if i == 2 else NAVY
    box(s, x, sy, sw2, 1.5, fc=bg2, rr=True)
    txt(s, title, x, sy + 0.12, sw2, 0.6,
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, x, sy + 0.72, sw2, 0.7,
        sz=8.5, color=RGBColor(0xBF,0xCA,0xE0), align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, ">", x + sw2 + 0.04, sy + 0.56, 0.2, 0.3,
            sz=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Haversine formula
box(s, 0.3, 2.78, 6.1, 0.82, fc=DARK, rr=True)
txt(s, "d = 2R arcsin( sqrt( sin²(Dlat/2) + cos(lat1)·cos(lat2)·sin²(Dlng/2) ) )",
    0.48, 2.9, 5.82, 0.42, sz=10.5, color=RGBColor(0x86,0xEF,0xAC))
txt(s, "R = 6371 km", 0.48, 3.3, 3, 0.24, sz=9, color=MUTED)

# Cost formula
box(s, 0.3, 3.74, 6.1, 1.28, fc=DARK, rr=True)
txt(s, "Cost (RWF)  =", 0.48, 3.84, 5.8, 0.38,
    sz=11, bold=True, color=ORANGE)
for li, line in enumerate(["2,000  (base fare)",
                            "+ 500  x  weight in tons",
                            "+ 200  x  distance in km"]):
    txt(s, line, 0.54, 4.22 + li * 0.28, 5.7, 0.28,
        sz=11, color=RGBColor(0xFD,0xE6,0x8A))

# vertical divider
box(s, 6.65, 1.05, 0.04, 5.5, fc=BORDER)

# ── RIGHT: mini ERD ───────────────────────────────────────────────────────────
tables = [
    ("USERS",     [("PK","user_id"), ("","email"), ("","role ENUM")]),
    ("CUSTOMERS", [("PK","customer_id"), ("FK","user_id")]),
    ("DRIVERS",   [("PK","driver_id"), ("FK","user_id"),
                   ("","is_verified"), ("","lat / lng")]),
    ("TRUCKS",    [("PK","truck_id"), ("FK","driver_id"), ("","capacity")]),
    ("BOOKINGS",  [("PK","booking_id"), ("FK","customer_id"),
                   ("FK","driver_id"), ("","status ENUM"), ("","cost")]),
    ("PAYMENTS",  [("PK","payment_id"), ("FK","booking_id UNIQUE"),
                   ("","status ENUM")]),
]
positions = [(6.9, 1.05), (6.9, 3.18), (9.05, 1.05),
             (9.05, 3.62), (10.88, 1.05), (10.88, 4.12)]
tw = 2.22; rh2 = 0.28

for (tname, tcols), (tx, ty) in zip(tables, positions):
    th = 0.34 + len(tcols) * rh2
    box(s, tx, ty, tw, th, fc=DARK, lc=ORANGE, lw=0.6, rr=True)
    box(s, tx, ty, tw, 0.34, fc=ORANGE, rr=True)
    txt(s, tname, tx + 0.06, ty + 0.05, tw - 0.12, 0.26,
        sz=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, (rtype, rname) in enumerate(tcols):
        ry2 = ty + 0.34 + ri * rh2
        tc = RGBColor(0xFD,0xE6,0x8A) if rtype=="PK" else \
             RGBColor(0x93,0xC5,0xFD) if rtype=="FK" else \
             RGBColor(0x86,0xEF,0xAC)
        if rtype:
            box(s, tx+0.05, ry2+0.04, 0.24, 0.2, fc=tc, rr=True)
            txt(s, rtype, tx+0.05, ry2+0.04, 0.24, 0.2,
                sz=6, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        txt(s, rname, tx+0.32, ry2+0.05, tw-0.38, 0.22,
            sz=7.5, color=RGBColor(0xE2,0xE8,0xF0))

# legend
lx = 6.9; ly = H - 0.55
for label, lc2 in [("PK", RGBColor(0xFD,0xE6,0x8A)),
                    ("FK", RGBColor(0x93,0xC5,0xFD)),
                    ("Col", RGBColor(0x86,0xEF,0xAC))]:
    box(s, lx, ly, 0.22, 0.18, fc=lc2, rr=True)
    txt(s, label, lx+0.25, ly, 0.5, 0.22, sz=7.5, color=MUTED)
    lx += 0.9

ftr(s, "Algorithm and Database")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 06 – SYSTEM ARCHITECTURE
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "System Architecture", 6)

layers = [
    ("CLIENT",      BLUE,                      BLUE_BG,
     "Browser   |   HTML / CSS / JS   |   Fetch API   |   Leaflet.js   |   Socket.IO Client",
     1.06, 0.82),
    ("APPLICATION", NAVY,                      BG,
     "Flask App Factory   |   auth   |   booking   |   driver   |   customer   |   payment   |   admin   |   tracking (SocketIO + eventlet)",
     2.02, 0.82),
    ("DATA",        RGBColor(0x06,0x4E,0x3B),  GREEN_BG,
     "PostgreSQL   |   SQLAlchemy ORM   |   8 tables   |   ENUM status fields   |   Flask-Migrate",
     2.98, 0.82),
    ("EXTERNAL",    RGBColor(0x78,0x35,0x00),  AMBER_BG,
     "MTN MoMo Sandbox API   |   /payment/callback webhook   |   OpenStreetMap tile server",
     3.94, 0.82),
]

proto_labels = ["HTTP / WebSocket (WS)", "SQLAlchemy ORM", "HTTPS  (requests lib)"]

for i, (label, accent, bg2, content, ty, th) in enumerate(layers):
    box(s, 0.3, ty, 0.22, th, fc=accent)
    box(s, 0.54, ty, W - 0.84, th, fc=bg2, lc=BORDER, lw=0.3)
    txt(s, label, 0.64, ty + 0.24, 1.9, 0.38, sz=10, bold=True, color=accent)
    txt(s, content, 2.6, ty + 0.24, W - 3.0, 0.38, sz=11, color=TEXT)
    if i < 3:
        txt(s, proto_labels[i], 1.0, ty + th + 0.02, 4.5, 0.2,
            sz=8, color=MUTED)

# Key interactions callout
box(s, 0.3, 4.94, W - 0.6, 0.72, fc=DARK, rr=True)
items = ["SocketIO rooms  ->  live GPS per booking",
         "MoMo callback  ->  async payment update",
         "bcrypt + Flask-Login  ->  secure sessions"]
for i, item in enumerate(items):
    x2 = 0.55 + i * 4.24
    circ(s, x2, 5.06, 0.34, ORANGE)
    txt(s, str(i+1), x2, 5.08, 0.34, 0.28,
        sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, item, x2 + 0.42, 5.1, 3.72, 0.36, sz=10, color=MUTED)

ftr(s, "System Architecture")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 07 – CODE QUALITY + TESTING
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "Code Quality and Testing", 7)

# ── LEFT: code block ──────────────────────────────────────────────────────────
box(s, 0.3, 1.05, 6.4, 5.15, fc=DARK, rr=True)
txt(s, "Role Decorator  (applied to all 44 routes)",
    0.48, 1.14, 6.1, 0.3, sz=9, bold=True, color=ORANGE)
code = [
    ("def require_admin(f):",                    RGBColor(0xFD,0xE6,0x8A)),
    ('    """Restrict route to admin."""',        RGBColor(0x6B,0x8C,0xBA)),
    ("    @wraps(f)",                            RGBColor(0x93,0xC5,0xFD)),
    ("    def decorated(*args, **kw):",          RGBColor(0xFD,0xE6,0x8A)),
    ('        if current_user.role != "admin":', RGBColor(0xE2,0xE8,0xF0)),
    ('            return error(403)',            RGBColor(0x86,0xEF,0xAC)),
    ("        return f(*args, **kw)",           RGBColor(0x86,0xEF,0xAC)),
    ("    return decorated",                     RGBColor(0xE2,0xE8,0xF0)),
    ("", WHITE),
    ("@admin_bp.route('/stats')",               RGBColor(0x93,0xC5,0xFD)),
    ("@login_required",                         RGBColor(0x93,0xC5,0xFD)),
    ("@require_admin          # <-- enforced",  ORANGE),
    ("def stats(): ...",                        RGBColor(0xFD,0xE6,0x8A)),
]
for li, (line, lc) in enumerate(code):
    txt(s, line, 0.46, 1.48 + li * 0.33, 6.1, 0.32, sz=10, color=lc)

# Standards pills
standards = ["PEP 8", "Docstrings", "DRY helpers", "Env secrets", "JSON envelope"]
sx = 0.36; sy2 = 5.72
for st in standards:
    sw3 = len(st)*0.085 + 0.28
    box(s, sx, sy2, sw3, 0.3, fc=NAVY2, rr=True)
    txt(s, st, sx+0.07, sy2+0.04, sw3-0.1, 0.24,
        sz=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx += sw3 + 0.12

# ── RIGHT: test results ───────────────────────────────────────────────────────
tests = [
    ("User Authentication",  "Pass"),
    ("Role Redirect",        "Pass"),
    ("Driver Search",        "Pass"),
    ("Booking Lifecycle",    "Pass"),
    ("Live GPS Tracking",    "Pass"),
    ("MoMo Payment (202)",   "Pass"),
    ("Admin Verification",   "Pass"),
]

box(s, 6.95, 1.05, 6.05, 0.38, fc=NAVY)
txt(s, "Manual Test", 7.05, 1.07, 3.4, 0.32, sz=9, bold=True, color=WHITE)
txt(s, "Result",      10.5, 1.07, 2.4, 0.32, sz=9, bold=True, color=WHITE)

for ri, (flow, result) in enumerate(tests):
    ry = 1.45 + ri * 0.68
    rc = BG if ri % 2 == 0 else WHITE
    box(s, 6.95, ry, 6.05, 0.66, fc=rc, lc=BORDER, lw=0.3)
    txt(s, flow, 7.05, ry + 0.14, 3.4, 0.38, sz=11, color=TEXT)
    box(s, 10.5, ry + 0.1, 2.38, 0.46, fc=GREEN_BG,
        lc=RGBColor(0x6E,0xE7,0xB7), lw=0.4, rr=True)
    circ(s, 10.6, ry + 0.12, 0.38, fc=GREEN)
    txt(s, "ok", 10.6, ry + 0.14, 0.38, 0.3,
        sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, result, 11.08, ry + 0.14, 1.7, 0.38,
        sz=11, bold=True, color=GREEN)

box(s, 6.95, 6.3, 6.05, 0.42, fc=DARK, rr=True)
txt(s, "Next: automated unit + integration tests  |  Flask test client  |  fixtures",
    7.08, 6.35, 5.85, 0.34, sz=9, color=MUTED)

ftr(s, "Code Quality and Testing")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 08 – TECHNICAL CHALLENGES
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "Technical Challenges and Solutions", 8)

challenges = [
    ("Real-Time Tracking",
     "HTTP polling overloads the server",
     "SocketIO room per booking  ->  sub-second GPS broadcast"),
    ("Driver Matching",
     "Naive filter loads all drivers into memory",
     "Single SQLAlchemy query: subquery + JOIN + boolean flags"),
    ("Async MoMo Payment",
     "Confirmation arrives via webhook, not immediately",
     "Pending record at initiation  ->  /callback updates status"),
    ("Role Enforcement",
     "Flask has no native RBAC",
     "@require_admin / @require_driver on every sensitive route"),
]

for i, (title, problem, solution) in enumerate(challenges):
    row = i // 2; col = i % 2
    x = 0.3 + col * 6.54; y = 1.06 + row * 3.1
    box(s, x, y, 6.28, 2.9, fc=BG, lc=BORDER, lw=0.4, rr=True)
    # top accent strip
    box(s, x, y, 6.28, 0.56, fc=NAVY, rr=True)
    num_circ(s, x + 0.14, y + 0.06, 0.44, i + 1)
    txt(s, title, x + 0.7, y + 0.1, 5.42, 0.38,
        sz=14, bold=True, color=WHITE)
    # problem
    box(s, x + 0.16, y + 0.72, 5.96, 0.82,
        fc=RED_BG, lc=RGBColor(0xFC,0xC2,0xD0), lw=0.4, rr=True)
    txt(s, "Problem", x + 0.28, y + 0.76, 1.1, 0.26,
        sz=8, bold=True, color=RED)
    txt(s, problem,  x + 0.28, y + 1.02, 5.68, 0.46,
        sz=10.5, color=RED)
    # solution
    box(s, x + 0.16, y + 1.66, 5.96, 1.08,
        fc=GREEN_BG, lc=RGBColor(0x6E,0xE7,0xB7), lw=0.4, rr=True)
    txt(s, "Solution", x + 0.28, y + 1.7, 1.1, 0.26,
        sz=8, bold=True, color=GREEN)
    txt(s, solution, x + 0.28, y + 1.96, 5.68, 0.7,
        sz=10.5, color=GREEN)

ftr(s, "Challenges and Solutions")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 09 – FEEDBACK + NEXT STEPS
# ═════════════════════════════════════════════════════════════════════════════
s = ns()
hdr(s, "Feedback Integration and Next Steps", 9)

# ── LEFT: before / after ──────────────────────────────────────────────────────
box(s, 0.3, 1.05, 6.3, 0.38, fc=NAVY)
txt(s, "BEFORE", 0.44, 1.08, 2.6, 0.3, sz=9, bold=True, color=RED_BG)
txt(s, "AFTER",  3.3,  1.08, 3.1, 0.3, sz=9, bold=True, color=GREEN_BG)

pairs = [
    ("DB queries inside route handlers",   "Extracted helper functions"),
    ("Inconsistent JSON responses",        "Shared success() / error() wrappers"),
    ("No RBAC on admin routes",            "@require_admin decorator added"),
    ("Status stored as VARCHAR",           "PostgreSQL ENUM types enforced"),
    ("Credentials in source code",         "All secrets moved to .env"),
]
for i, (before, after) in enumerate(pairs):
    ry = 1.45 + i * 1.0
    # before
    box(s, 0.3, ry, 2.86, 0.86, fc=RED_BG, lc=RGBColor(0xFC,0xC2,0xD0), lw=0.4, rr=True)
    circ(s, 0.4, ry + 0.2, 0.4, fc=RED)
    txt(s, "x", 0.4, ry + 0.22, 0.4, 0.3,
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, before, 0.9, ry + 0.2, 2.18, 0.5, sz=9.5, color=RED, wrap=True)
    # arrow
    txt(s, ">", 3.2, ry + 0.3, 0.26, 0.28, sz=13, bold=True, color=MUTED)
    # after
    box(s, 3.5, ry, 3.1, 0.86, fc=GREEN_BG, lc=RGBColor(0x6E,0xE7,0xB7), lw=0.4, rr=True)
    circ(s, 3.6, ry + 0.2, 0.4, fc=GREEN)
    txt(s, "ok", 3.6, ry + 0.22, 0.4, 0.3,
        sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, after, 4.1, ry + 0.2, 2.38, 0.5, sz=9.5, color=GREEN, wrap=True)

# vertical divider
box(s, 6.82, 1.0, 0.04, 5.65, fc=BORDER)

# ── RIGHT: next steps timeline ────────────────────────────────────────────────
txt(s, "NEXT TECHNICAL MILESTONES", 7.0, 1.05, 6.1, 0.28,
    sz=9, bold=True, color=ORANGE)

milestones = [
    ("Automated Tests",     "[Member 1]", "Unit + integration, Flask test client"),
    ("Email Notifications", "[Member 2]", "Flask-Mail / SendGrid SMTP"),
    ("MoMo Production",     "[Member 2]", "Migrate sandbox credentials + idempotency"),
    ("Document Upload",     "[Member 3]", "Driver licence and vehicle image upload"),
    ("PWA Support",         "[Member 4]", "Service worker + offline cache"),
]

# Timeline bar
box(s, 7.2, 1.72, 0.06, 5.38, fc=BORDER)   # vertical rail

for i, (title, member, desc) in enumerate(milestones):
    cy = 1.62 + i * 1.06
    circ(s, 7.0, cy, 0.46, fc=ORANGE)
    txt(s, str(i+1), 7.0, cy + 0.06, 0.46, 0.36,
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 7.6, cy, 5.68, 0.9, fc=BG, lc=BORDER, lw=0.3, rr=True)
    txt(s, title,  7.72, cy + 0.06, 4.0, 0.34, sz=11, bold=True, color=NAVY)
    txt(s, member, 11.2, cy + 0.06, 2.0, 0.28, sz=8.5, bold=True,
        color=ORANGE, align=PP_ALIGN.RIGHT)
    txt(s, desc,   7.72, cy + 0.42, 5.42, 0.38, sz=9.5, color=MUTED)

ftr(s, "Feedback and Next Steps")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – THANK YOU + REFERENCES
# ═════════════════════════════════════════════════════════════════════════════
s = ns(NAVY)
box(s, 0, 0,        W, 0.16, fc=ORANGE)
box(s, 0, H - 0.16, W, 0.16, fc=ORANGE)

# Left closing block
box(s, 0.65, 0.6, 0.08, 4.4, fc=ORANGE)
txt(s, "Thank You", 0.95, 0.55, 8, 1.6,
    sz=60, bold=True, color=WHITE)
txt(s, "Questions and Discussion",
    0.95, 2.12, 8, 0.55, sz=18,
    color=RGBColor(0xBF,0xCA,0xE0))
box(s, 0.95, 2.78, 2.8, 0.05, fc=ORANGE)

links = [("GitHub", "[ INSERT LINK ]"),
         ("Figma",  "[ INSERT LINK ]")]
for i, (label, url) in enumerate(links):
    box(s, 0.95, 2.98 + i*0.66, 0.88, 0.46, fc=ORANGE, rr=True)
    txt(s, label, 0.95, 3.02 + i*0.66, 0.88, 0.38,
        sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, url, 1.93, 3.05 + i*0.66, 5.5, 0.36,
        sz=10, color=RGBColor(0xBF,0xCA,0xE0))

# Stack pills
stack_labels = ["Flask", "PostgreSQL", "Flask-SocketIO",
                "Leaflet.js", "MTN MoMo API", "bcrypt"]
px = 0.95; py = 4.5
for lbl in stack_labels:
    pw2 = len(lbl)*0.085 + 0.3
    box(s, px, py, pw2, 0.32, fc=NAVY2, rr=True)
    txt(s, lbl, px+0.08, py+0.04, pw2-0.12, 0.26,
        sz=9, color=WHITE, align=PP_ALIGN.CENTER)
    px += pw2 + 0.12
    if px > 6.5:
        px = 0.95; py += 0.42

# Right: mini references
box(s, 7.4, 0.4, 5.6, 6.64, fc=NAVY2, rr=True)
txt(s, "REFERENCES", 7.58, 0.52, 5.28, 0.3,
    sz=9, bold=True, color=ORANGE)
refs = [
    ("Flask 3.x",          "flask.palletsprojects.com"),
    ("Flask-SQLAlchemy",   "flask-sqlalchemy.palletsprojects.com"),
    ("Flask-SocketIO",     "flask-socketio.readthedocs.io"),
    ("Flask-Login",        "flask-login.readthedocs.io"),
    ("psycopg2-binary",    "psycopg.org"),
    ("python-dotenv",      "pypi.org/project/python-dotenv"),
    ("Leaflet.js",         "leafletjs.com"),
    ("Socket.IO Client",   "socket.io"),
    ("MTN MoMo API",       "momodeveloper.mtn.com"),
    ("OpenStreetMap",      "openstreetmap.org"),
    ("Haversine Formula",  "Sinnott, Sky & Telescope, 1984"),
    ("App Factory Pattern","flask.palletsprojects.com/patterns"),
]
for ri, (name, url) in enumerate(refs):
    ry = 0.92 + ri * 0.47
    rc = NAVY if ri % 2 == 0 else NAVY2
    box(s, 7.44, ry, 5.52, 0.45, fc=rc)
    txt(s, name, 7.54, ry + 0.08, 2.4, 0.3, sz=9, bold=True, color=WHITE)
    txt(s, url,  9.98, ry + 0.08, 2.9, 0.3, sz=8.5, color=MUTED)


# ── SAVE ──────────────────────────────────────────────────────────────────────
out = "/Users/nkemachia/Documents/Trimester_4/CargoConnect/CargoConnect_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  |  {len(prs.slides)} slides")
